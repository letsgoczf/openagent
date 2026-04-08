"""从多种文档格式抽取纯文本，按「页/段」返回供入库与向量化。"""

from __future__ import annotations

import csv
import email
import io
import json
import zipfile
from email.policy import default as email_policy_default
from html.parser import HTMLParser
from pathlib import Path
from xml.etree import ElementTree as ET

import pdfplumber
import xlrd
import yaml
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from striprtf.striprtf import rtf_to_text

# 单段过长时按字符切分为多「页」，避免单次 embedding 过大
_MAX_SEGMENT_CHARS = 12_000

# 常见 OLE Compound Document（旧版 .doc）
_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


class DocumentExtractionError(Exception):
    """当前格式无法解析或缺少依赖时抛出。"""


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        t = data.strip()
        if t:
            self._parts.append(t)

    def text(self) -> str:
        return "\n".join(self._parts)


def _html_to_text(raw: str) -> str:
    p = _HTMLTextExtractor()
    p.feed(raw)
    p.close()
    return p.text()


def _segment_long_text(text: str) -> list[str]:
    text = (text or "").strip()
    if not text:
        return []
    if len(text) <= _MAX_SEGMENT_CHARS:
        return [text]
    chunks: list[str] = []
    rest = text
    while rest:
        if len(rest) <= _MAX_SEGMENT_CHARS:
            chunks.append(rest.strip())
            break
        cut = rest.rfind("\n\n", 0, _MAX_SEGMENT_CHARS)
        if cut < _MAX_SEGMENT_CHARS // 4:
            cut = rest.rfind("\n", 0, _MAX_SEGMENT_CHARS)
        if cut < _MAX_SEGMENT_CHARS // 4:
            cut = _MAX_SEGMENT_CHARS
        piece = rest[:cut].strip()
        if piece:
            chunks.append(piece)
        rest = rest[cut:].lstrip()
    return [c for c in chunks if c]


def _decode_plain(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp936", "gbk", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _looks_binary(sample: bytes) -> bool:
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    printable = sum(1 for b in sample if 32 <= b < 127 or b in (9, 10, 13))
    return printable / len(sample) < 0.55


def _extract_pdf_pages(file_bytes: bytes) -> list[str]:
    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for p in pdf.pages:
            txt = (p.extract_text() or "").strip()
            pages.append(txt)
    return pages


def _extract_docx(file_bytes: bytes) -> list[str]:
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append(t)
    for tbl in doc.tables:
        rows_txt: list[str] = []
        for row in tbl.rows:
            cells = [" ".join((c.text or "").split()) for c in row.cells]
            rows_txt.append("\t".join(cells))
        if rows_txt:
            parts.append("\n".join(rows_txt))
    return _segment_long_text("\n\n".join(parts))


def _pptx_walk_shapes(shapes: object, chunk: list[str]) -> None:
    for shape in shapes:  # type: ignore[assignment]
        st = shape.shape_type  # type: ignore[attr-defined]
        if st == MSO_SHAPE_TYPE.TABLE:
            tbl = shape.table  # type: ignore[attr-defined]
            for row in tbl.rows:
                cells = [" ".join((c.text or "").split()) for c in row.cells]
                chunk.append("\t".join(cells))
        elif st == MSO_SHAPE_TYPE.GROUP:
            _pptx_walk_shapes(shape.shapes, chunk)  # type: ignore[attr-defined]
        elif getattr(shape, "has_text_frame", False) and shape.text:  # type: ignore[attr-defined]
            t = " ".join(shape.text.split())
            if t:
                chunk.append(t)


def _extract_pptx(file_bytes: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []
    for slide in prs.slides:
        chunk: list[str] = []
        _pptx_walk_shapes(slide.shapes, chunk)
        if chunk:
            slides.append("\n".join(chunk))
    merged = "\n\n".join(slides)
    return _segment_long_text(merged) if merged.strip() else []


def _extract_xlsx(file_bytes: bytes) -> list[str]:
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in wb:
            rows_txt: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = []
                for c in row:
                    if c is None:
                        cells.append("")
                    else:
                        cells.append(str(c).strip())
                if any(cells):
                    rows_txt.append("\t".join(cells))
            if rows_txt:
                parts.append(f"## {sheet.title}\n" + "\n".join(rows_txt))
    finally:
        wb.close()
    return _segment_long_text("\n\n".join(parts))


def _extract_xls(file_bytes: bytes) -> list[str]:
    book = xlrd.open_workbook(file_contents=file_bytes)
    parts: list[str] = []
    for si in range(book.nsheets):
        sheet = book.sheet_by_index(si)
        rows_txt: list[str] = []
        for ri in range(sheet.nrows):
            cells = [str(sheet.cell_value(ri, ci)).strip() for ci in range(sheet.ncols)]
            if any(cells):
                rows_txt.append("\t".join(cells))
        if rows_txt:
            parts.append(f"## {sheet.name}\n" + "\n".join(rows_txt))
    return _segment_long_text("\n\n".join(parts))


def _extract_rtf(file_bytes: bytes) -> list[str]:
    raw = file_bytes.decode("latin-1", errors="replace")
    text = rtf_to_text(raw)
    return _segment_long_text(text)


def _extract_odt(file_bytes: bytes) -> list[str]:
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            xml_data = z.read("content.xml")
    except (KeyError, zipfile.BadZipFile) as e:
        raise DocumentExtractionError("无法读取 ODT（需为有效 OpenDocument 压缩包）") from e
    root = ET.fromstring(xml_data)
    full = "".join(root.itertext())
    return _segment_long_text(full)


def _extract_epub(file_bytes: bytes) -> list[str]:
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            names = z.namelist()
            html_names = sorted(
                n
                for n in names
                if n.lower().endswith((".xhtml", ".html", ".htm"))
                and not n.startswith("__MACOSX")
            )
            blobs: list[str] = []
            for name in html_names:
                try:
                    raw = z.read(name)
                    blobs.append(_html_to_text(raw.decode("utf-8", errors="ignore")))
                except KeyError:
                    pass
    except zipfile.BadZipFile as e:
        raise DocumentExtractionError("无法读取 EPUB（需为有效 zip）") from e
    merged = "\n\n".join(b for b in blobs if b.strip())
    return _segment_long_text(merged) if merged.strip() else []


def _extract_eml(file_bytes: bytes) -> list[str]:
    msg = email.message_from_bytes(file_bytes, policy=email_policy_default)
    subj = (msg.get("subject") or "").strip()
    lines: list[str] = [subj] if subj else []
    plain_parts: list[str] = []
    html_parts: list[str] = []

    for part in msg.walk():
        ctype = part.get_content_type()
        if ctype == "text/plain":
            try:
                payload = part.get_content()
            except Exception:  # noqa: BLE001
                continue
            if isinstance(payload, str) and payload.strip():
                plain_parts.append(payload.strip())
        elif ctype == "text/html":
            try:
                raw = part.get_payload(decode=True) or b""
                html_parts.append(_html_to_text(raw.decode("utf-8", errors="ignore")))
            except Exception:  # noqa: BLE001
                continue

    if plain_parts:
        lines.extend(plain_parts)
    elif html_parts:
        lines.extend(html_parts)
    return _segment_long_text("\n\n".join(lines))


def _extract_csv(file_bytes: bytes) -> list[str]:
    text = _decode_plain(file_bytes)
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(row) for row in reader]
    return _segment_long_text("\n".join(rows))


def _extract_json_bytes(file_bytes: bytes) -> list[str]:
    text = _decode_plain(file_bytes)
    try:
        obj = json.loads(text)
    except json.JSONDecodeError as e:
        raise DocumentExtractionError(f"不是合法 JSON：{e}") from e
    pretty = json.dumps(obj, ensure_ascii=False, indent=2)
    return _segment_long_text(pretty)


def _extract_xml_bytes(file_bytes: bytes) -> list[str]:
    try:
        root = ET.fromstring(file_bytes)
    except ET.ParseError as e:
        raise DocumentExtractionError(f"XML 解析失败：{e}") from e
    full = "".join(root.itertext())
    return _segment_long_text(full)


def _extract_jsonl_bytes(file_bytes: bytes) -> list[str]:
    out: list[str] = []
    for line in _decode_plain(file_bytes).splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            out.append(json.dumps(json.loads(line), ensure_ascii=False))
        except json.JSONDecodeError:
            continue
    return _segment_long_text("\n".join(out))


def _extract_yaml_bytes(file_bytes: bytes) -> list[str]:
    text = _decode_plain(file_bytes)
    try:
        obj = yaml.safe_load(text)
    except yaml.YAMLError as e:
        raise DocumentExtractionError(f"不是合法 YAML：{e}") from e
    if obj is None:
        return []
    pretty = yaml.safe_dump(obj, allow_unicode=True, default_flow_style=False)
    return _segment_long_text(pretty)


def _sniff_ooxml_kind(file_bytes: bytes) -> str | None:
    if not file_bytes.startswith(b"PK"):
        return None
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            names = set(z.namelist())
    except zipfile.BadZipFile:
        return None
    if "word/document.xml" in names:
        return "docx"
    if "xl/workbook.xml" in names:
        return "xlsx"
    if "ppt/presentation.xml" in names:
        return "pptx"
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            if "mimetype" in z.namelist():
                mt = z.read("mimetype").decode("utf-8", errors="ignore").strip()
                if "opendocument.text" in mt and "content.xml" in names:
                    return "odt"
                if mt == "application/epub+zip":
                    return "epub"
    except (zipfile.BadZipFile, KeyError):
        pass
    if "content.xml" in names and "META-INF/manifest.xml" in names:
        return "odt"
    if "META-INF/container.xml" in names:
        return "epub"
    return None


def _resolve_handler_key(file_bytes: bytes, ext: str) -> str:
    ext_l = ext.lower()
    plain_exts = {
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".log",
        ".ini",
        ".cfg",
        ".toml",
        ".tsv",
        ".sql",
        ".sh",
        ".bat",
        ".c",
        ".h",
        ".cpp",
        ".hpp",
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".java",
        ".go",
        ".rs",
        ".rb",
        ".php",
        ".vue",
        ".css",
        ".scss",
        ".env",
    }
    mapping: dict[str, str] = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".pptx": "pptx",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".rtf": "rtf",
        ".odt": "odt",
        ".epub": "epub",
        ".html": "html",
        ".htm": "html",
        ".xhtml": "html",
        ".csv": "csv",
        ".json": "json",
        ".jsonl": "jsonl",
        ".xml": "xml",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".eml": "eml",
    }
    if ext_l in mapping:
        return mapping[ext_l]
    if ext_l in plain_exts:
        return "plain"

    if file_bytes.startswith(b"%PDF-"):
        return "pdf"
    if file_bytes.startswith(_OLE_MAGIC):
        return "ole"
    sniffed = _sniff_ooxml_kind(file_bytes)
    if sniffed:
        return sniffed

    sample = file_bytes[:8192]
    if not _looks_binary(sample):
        return "plain"

    raise DocumentExtractionError(
        "无法识别该二进制文件类型，请使用常见文档后缀（如 .pdf/.docx/.txt）或转换为受支持格式。"
    )


def extract_document_pages(file_bytes: bytes, filename: str) -> list[str]:
    """
    按文档内容返回若干段文本（PDF 按页；长文按段切分），每段将对应一个 chunk。
    """
    if not file_bytes:
        raise DocumentExtractionError("空文件")

    path = Path(filename or "upload")
    ext = path.suffix
    key = _resolve_handler_key(file_bytes, ext)

    if key == "ole":
        if ext.lower() == ".doc" or file_bytes.startswith(_OLE_MAGIC):
            raise DocumentExtractionError(
                "不支持旧版 Word .doc（OLE），请另存为 .docx 或导出为 PDF 后再导入。"
            )
        raise DocumentExtractionError("不支持的 OLE/二进制文档格式。")

    handlers: dict[str, object] = {
        "pdf": lambda: _extract_pdf_pages(file_bytes),
        "docx": lambda: _extract_docx(file_bytes),
        "pptx": lambda: _extract_pptx(file_bytes),
        "xlsx": lambda: _extract_xlsx(file_bytes),
        "xls": lambda: _extract_xls(file_bytes),
        "rtf": lambda: _extract_rtf(file_bytes),
        "odt": lambda: _extract_odt(file_bytes),
        "epub": lambda: _extract_epub(file_bytes),
        "html": lambda: _segment_long_text(_html_to_text(_decode_plain(file_bytes))),
        "csv": lambda: _extract_csv(file_bytes),
        "json": lambda: _extract_json_bytes(file_bytes),
        "jsonl": lambda: _extract_jsonl_bytes(file_bytes),
        "xml": lambda: _extract_xml_bytes(file_bytes),
        "yaml": lambda: _extract_yaml_bytes(file_bytes),
        "eml": lambda: _extract_eml(file_bytes),
        "plain": lambda: _segment_long_text(_decode_plain(file_bytes)),
    }

    fn = handlers.get(key)
    if not fn:
        raise DocumentExtractionError(f"内部错误：未知处理器 {key}")

    try:
        pages = fn()  # type: ignore[misc,operator]
    except DocumentExtractionError:
        raise
    except Exception as e:  # noqa: BLE001
        raise DocumentExtractionError(f"文档解析失败：{e}") from e
    if not isinstance(pages, list):
        raise DocumentExtractionError("抽取结果异常")
    # 统一去掉全空段
    return [p for p in pages if isinstance(p, str) and p.strip()]
